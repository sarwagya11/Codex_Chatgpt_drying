"""v2: start from refined deck, sharpen slides 8 & 11, drop slides 12-14, fix page numbers."""
from pathlib import Path
from pptx import Presentation

SRC = Path(r"C:\Users\sarwa\Downloads\E2_vs_E3_refined.pptx")
DST = Path(r"C:\Users\sarwa\Downloads\E2_vs_E3_final.pptx")

p = Presentation(SRC)


def replace_in_slide(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                for old, new in mapping.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


# --- Slide 8: sharpen winners table notes ---
s8 = p.slides[7]
replace_in_slide(s8, {
    "Best energy economy overall.":
        "Best SEC; mean -7.8% vs E3 (no-VPD), -8.8% (VPD-on).",
    "Higher COP does not translate to SEC.":
        "Higher COP, partial-lift control. SEC penalty +7.8%/+8.8%.",
    "Cold evap forces high Q_cond — energetically expensive.":
        "Fastest t_h. SEC ~ +8.6% vs E2 (no-VPD), +3.5% (VPD-on).",
    "COP ≠ SEC predictor.": "COP is not a SEC predictor.",
    "Right KPI — SEC, not COP.": "Rank by SEC and SMER, not COP.",
})

# --- Slide 11: sharpen audit ---
s11 = p.slides[10]
replace_in_slide(s11, {
    "92 canonical runs screened across E1 / E2 / E3 × locations × seasons × VPD on/off.":
        "92 canonical runs (E1/E2/E3 × 4 locations × 4 seasons × VPD on/off), 7 anomaly classes screened.",
    "0 across all runs":
        "0 / 92",
    "E2 / BTN / autumn / VPD reaches 99 % moisture target at 72 h cutoff.":
        "E2/BTN/autumn/VPD: 99% of moisture target at 72 h cutoff (humid lowland stress case).",
    "E3 partial-lift transients with COP > 8 — physically valid.":
        "E3 partial-lift transients, COP > 8 at low T_cond targets (Carnot-valid).",
})

# --- Drop slides 12, 13, 14 ---
xml_slides = p.slides._sldIdLst
sldIds = list(xml_slides)
# remove relationships then nodes
for sldId in sldIds[11:]:
    rId = sldId.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    p.part.drop_rel(rId)
    xml_slides.remove(sldId)

# --- Fix page numbers "NN / 14" -> "NN / 11" on remaining 11 slides ---
for slide in p.slides:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if "/ 14" in run.text:
                    run.text = run.text.replace("/ 14", "/ 11")

p.save(DST)
print(f"Saved {DST}; slide count = {len(p.slides)}")

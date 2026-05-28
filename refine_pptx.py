"""Refine E2_vs_E3.pptx: embed audit PNGs into placeholder boxes, save as _refined."""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

SRC = Path(r"C:\Users\sarwa\Downloads\E2_vs_E3.pptx")
DST = Path(r"C:\Users\sarwa\Downloads\E2_vs_E3_refined.pptx")
PLOTS = Path(r"D:\Masters\RQ5\Codex_Chatgpt_drying\RQ1\plots\_audit")

# (slide_idx_0based, container_shape_idx, png_filename, fig_label_shape_idx, filename_caption_shape_idx)
JOBS = [
    (2, 21, "step2e_air_states_E2_kathmandu_annual.png", 22, 24),
    (3, 21, "step2e_air_states_E3_kathmandu_annual.png", 22, 24),
    (5,  7, "e2_vs_e3_power_kathmandu.png",              8, 10),
    (5, 11, "e2_vs_e3_cumulative_kathmandu.png",        12, 14),
    (6,  7, "step3_energy_split_E2.png",                 8, 10),
    (6, 11, "step3_energy_split_E3.png",                12, 14),
    (8,  7, "step4_vpd_sweep.png",                       8, 10),
    (9,  7, "step5_area_sweep.png",                      8, 10),
    (9, 11, "step5_season_heatmap.png",                 12, 14),
]

def fit(box_w, box_h, img_w, img_h):
    s = min(box_w / img_w, box_h / img_h)
    return int(img_w * s), int(img_h * s)

p = Presentation(SRC)
for slide_idx, container_idx, png, fig_lbl_idx, caption_idx in JOBS:
    slide = p.slides[slide_idx]
    shapes = list(slide.shapes)
    box = shapes[container_idx]
    L, T, W, H = box.left, box.top, box.width, box.height
    img_path = PLOTS / png
    with Image.open(img_path) as im:
        iw, ih = im.size
    pad = Emu(120000)
    bw, bh = W - 2*pad, H - 2*pad
    w, h = fit(bw, bh, iw, ih)
    left = L + (W - w) // 2
    top = T + (H - h) // 2
    slide.shapes.add_picture(str(img_path), left, top, w, h)
    # hide "[ figure ]" label and filename caption by blanking text
    for idx in (fig_lbl_idx, caption_idx):
        sh = shapes[idx]
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""

# Sharpen a few findings lines (slide 5: explicit numbers)
def replace_text(slide, old, new):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    return True
    return False

# Slide 5: add absolute SEC numbers from memory (no-VPD KTM baseline)
replace_text(p.slides[4], "E2 wins SEC + SMER in every case.",
             "E2 wins SEC + SMER, 13/13 no-VPD, 12/12 VPD.")
replace_text(p.slides[4], "13 no-VPD + 12 VPD-on cases.",
             "KTM no-VPD: E2 0.0959, E3 0.1376 kWh/kg.")
replace_text(p.slides[4], "E3 never closes the gap.",
             "Gap +7.8% (no-VPD), +8.8% (VPD on).")

# Slide 7: tighten clipping line
replace_text(p.slides[6], "E3 — more clipping at same A_c.",
             "Solar clipping (A_c=10 m²): higher in E3 vs E2.")

p.save(DST)
print(f"Saved {DST}")

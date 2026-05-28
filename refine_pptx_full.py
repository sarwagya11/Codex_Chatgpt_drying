"""End-to-end rebuild: from original deck, embed PNGs (incl. fresh heatmap),
sharpen slides 5/7/8/11, drop slides 12-14, fix page numbers."""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

SRC = Path(r"C:\Users\sarwa\Downloads\E2_vs_E3.pptx")
DST = Path(r"C:\Users\sarwa\Downloads\E2_vs_E3_final.pptx")
PLOTS = Path(r"D:\Masters\RQ5\Codex_Chatgpt_drying\RQ1\plots\_audit")

JOBS = [
    (2, 21, "step2e_air_states_E2_kathmandu_annual.png", 22, 24),
    (3, 21, "step2e_air_states_E3_kathmandu_annual.png", 22, 24),
    (5,  7, "e2_vs_e3_power_kathmandu.png",              8, 10),
    (5, 11, "e2_vs_e3_cumulative_kathmandu.png",        12, 14),
    (6,  7, "step3_energy_split_E2.png",                 8, 10),
    (6, 11, "step3_energy_split_E3.png",                12, 14),
    (8,  7, "step4_vpd_sweep.png",                       8, 10),
    (9,  7, "step5_area_sweep.png",                      8, 10),
    (9, 11, "step5_season_heatmap.png",                 12, 14),
]


def fit(bw, bh, iw, ih):
    s = min(bw / iw, bh / ih)
    return int(iw * s), int(ih * s)


def replace_paragraph(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            for old, new in mapping.items():
                if old in full:
                    runs = para.runs
                    if not runs:
                        continue
                    runs[0].text = full.replace(old, new)
                    for r in runs[1:]:
                        r.text = ""
                    break


p = Presentation(SRC)

# Embed figures
for slide_idx, container_idx, png, fig_lbl_idx, caption_idx in JOBS:
    slide = p.slides[slide_idx]
    shapes = list(slide.shapes)
    box = shapes[container_idx]
    L, T, W, H = box.left, box.top, box.width, box.height
    img_path = PLOTS / png
    with Image.open(img_path) as im:
        iw, ih = im.size
    pad = Emu(120000)
    w, h = fit(W - 2*pad, H - 2*pad, iw, ih)
    left = L + (W - w) // 2
    top = T + (H - h) // 2
    slide.shapes.add_picture(str(img_path), left, top, w, h)
    for idx in (fig_lbl_idx, caption_idx):
        sh = shapes[idx]
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""

# Slides 3 & 4 — rephrase closure line in plain English
for idx in (2, 3):
    replace_paragraph(p.slides[idx], {
        "Closure \u2014 Q_evap = Q_cond \u2212 \u03b7\u00b7W_comp.":
            "Evap is sized to absorb the heat the cond rejects, net of compressor work (1st law).",
    })

# Slide 5 — headline numbers
replace_paragraph(p.slides[4], {
    "E2 wins SEC + SMER in every case.":
        "E2 wins SEC + SMER, 13/13 no-VPD, 12/12 VPD.",
    "13 no-VPD + 12 VPD-on cases.":
        "KTM no-VPD: E2 0.0959, E3 0.1376 kWh/kg.",
    "E3 never closes the gap.":
        "Gap +7.8% (no-VPD), +8.8% (VPD on).",
})

# Slide 7 — clipping
replace_paragraph(p.slides[6], {
    "E3 \u2014 more clipping at same A_c.":
        "Solar clipping (A_c=10 m\u00b2): higher in E3 vs E2.",
})

# Slide 8 — winners table
replace_paragraph(p.slides[7], {
    "Best energy economy overall.":
        "Best SEC; mean -7.8% vs E3 (no-VPD), -8.8% (VPD-on).",
    "Higher COP does not translate to SEC.":
        "Higher COP, partial-lift control. SEC penalty +7.8%/+8.8%.",
    "Cold evap forces high Q_cond \u2014 energetically expensive.":
        "Fastest t_h. SEC ~ +8.6% vs E2 (no-VPD), +3.5% (VPD-on).",
    "COP \u2260 SEC predictor.": "COP is not a SEC predictor.",
    "Right KPI \u2014 SEC, not COP.": "Rank by SEC and SMER, not COP.",
})

# Slide 11 — audit (simplified, plain English)
replace_paragraph(p.slides[10], {
    "92 canonical runs screened across E1 / E2 / E3 \u00d7 locations \u00d7 seasons \u00d7 VPD on/off.":
        "Every run (92 in total, all configs and conditions) was checked for physical sanity.",
    "First-law violations":
        "Energy balance",
    "Capacity flags":
        "Equipment at limit",
    "T_chamber deficits":
        "Chamber reaches 45 \u00b0C",
    "Bypass churn events":
        "Bypass control stable",
    "0 across all runs": "0 / 92 runs flagged",
    "Warn ( \u00d71 )": "Minor warning (\u00d71)",
    "E2 / BTN / autumn / VPD reaches 99 % moisture target at 72 h cutoff.":
        "Biratnagar in autumn (E2 with VPD bypass) reaches 99 % dryness at the 72 h cutoff, just short of full target.",
    "Info ( \u00d72 )": "Info (\u00d72)",
    "E3 partial-lift transients with COP > 8 \u2014 physically valid.":
        "E3 momentarily shows COP > 8 when the heat pump runs at very low lift. Physically valid (high Carnot ceiling).",
})

# Drop slides 12-14
xml_slides = p.slides._sldIdLst
sldIds = list(xml_slides)
for sldId in sldIds[11:]:
    rId = sldId.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    p.part.drop_rel(rId)
    xml_slides.remove(sldId)

# Fix page numbers
for slide in p.slides:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if "/ 14" in run.text:
                    run.text = run.text.replace("/ 14", "/ 11")

p.save(DST)
print(f"Saved {DST}; slides = {len(p.slides)}")

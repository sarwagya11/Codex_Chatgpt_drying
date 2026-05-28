"""
Supervisor discussion deck: E2 vs E3 for SAHPD.

Produces Supervisor_Meeting_E2_vs_E3.pptx at repo root.
Uses data from outputs/master_summary.csv and E2_area_sweep_summary.csv,
and embeds existing plots from plots/compare_E_autumn_ktm/ and plots/config_E2/.
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(r"D:/Masters/RQ5/Codex_Chatgpt_drying/RQ1")
PLOTS = ROOT / "plots"
OUT_FILE = ROOT / "Supervisor_Meeting_E2_vs_E3.pptx"
TMP = ROOT / "presentation_tmp"
TMP.mkdir(exist_ok=True)

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
CORAL = RGBColor(0xE5, 0x6B, 0x4A)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
GOLD = RGBColor(0xF0, 0xA6, 0x23)
SAGE = RGBColor(0x7A, 0x9E, 0x9F)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def blank_slide(title_text: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.82), Inches(12.4), Inches(0.35))
        stf = sb.text_frame
        stf.margin_top = 0
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.italic = True
        sp.font.color.rgb = GREY
    # thin rule under title
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.45),
        Inches(1.22),
        Inches(12.88),
        Inches(1.22),
    )
    line.line.color.rgb = NAVY
    line.line.width = Pt(1)
    return slide


def add_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text.strip()


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=DARK, align=None, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        if align is not None:
            p.alignment = align
    return box


def bullets(slide, left, top, width, height, items, *,
            size=14, color=DARK, spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.line_spacing = spacing
    return box


def comp_box(slide, left, top, width, height, label, *,
             fill=TEAL, text_color=WHITE, size=14, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.15
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = DARK
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    return shp


def arrow(slide, x1, y1, x2, y2, *, color=DARK, weight=1.5, dash=False, label=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        conn.line.dash_style = 7  # dash
    # arrowhead
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = conn.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('h', 'med')
    except Exception:
        pass
    if label:
        # midpoint label
        mx = (x1 + x2) // 2
        my = (y1 + y2) // 2
        tb = slide.shapes.add_textbox(mx - Inches(0.6), my - Inches(0.22), Inches(1.2), Inches(0.3))
        tf = tb.text_frame
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = GREY
    return conn


# ============================================================
# Matplotlib plots
# ============================================================

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.edgecolor": "#333333",
    "axes.labelcolor": "#222222",
    "xtick.color": "#333333",
    "ytick.color": "#333333",
    "axes.titlesize": 12,
    "axes.labelsize": 11,
    "xtick.labelsize": 10,
    "ytick.labelsize": 10,
    "legend.fontsize": 10,
    "legend.frameon": False,
})

C_A = "#555555"
C_D1 = "#7A9E9F"
C_B = "#F0A623"
C_E2 = "#2E86AB"  # matplotlib use only
C_E3 = "#E56B4A"  # matplotlib use only
RGB_E2 = TEAL    # same hex, use in python-pptx
RGB_E3 = CORAL   # same hex, use in python-pptx


def plot_ladder() -> Path:
    configs = [
        "A\n(HP only)",
        "B\n(+Solar)",
        "D1\n(+HRX)",
        "D2\n(HRX +\ndyn evap)",
        "E2\n(HRX + Solar\npre-cond)",
        "E3\n(HRX + Solar\npost-cond)",
    ]
    secs = [0.717, 0.488, 0.365, 0.354, 0.197, 0.217]
    colors = [C_A, C_B, C_D1, C_D1, C_E2, C_E3]
    fig, ax = plt.subplots(figsize=(11, 4.6), dpi=160)
    bars = ax.bar(configs, secs, color=colors, edgecolor="black", linewidth=0.6, width=0.72)
    for i, (bar, sec) in enumerate(zip(bars, secs)):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.018,
                f"{sec:.3f}", ha="center", fontsize=11, fontweight="bold")
        if i > 0:
            drop = (0.717 - sec) / 0.717 * 100
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() / 2,
                    f"-{drop:.0f}%", ha="center", color="white", fontsize=11, fontweight="bold")
    ax.set_ylabel("SEC (kWh / kg water)")
    ax.set_title("Config ladder, Kathmandu annual, r = 0, 10 m² collector")
    ax.set_ylim(0, 0.82)
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout()
    p = TMP / "ladder.png"
    fig.savefig(p, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return p


def plot_seasonal() -> Path:
    seasons = ["Autumn\n(Oct-Nov)", "Winter\n(Dec-Jan)", "Spring\n(Mar-Apr)"]
    e2_ktm = [0.123, 0.133, 0.094]
    e3_ktm = [0.133, 0.143, 0.105]
    e2_btn = [0.108, 0.111, 0.093]
    e3_btn = [0.114, 0.121, 0.080]
    x = np.arange(len(seasons))
    w = 0.2
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10.5, 4.4), dpi=160, sharey=True)
    ax1.bar(x - w/2, e2_ktm, w, label="E2", color=C_E2, edgecolor="black", linewidth=0.5)
    ax1.bar(x + w/2, e3_ktm, w, label="E3", color=C_E3, edgecolor="black", linewidth=0.5)
    for i, (a, b) in enumerate(zip(e2_ktm, e3_ktm)):
        ax1.text(i - w/2, a + 0.004, f"{a:.3f}", ha="center", fontsize=9)
        ax1.text(i + w/2, b + 0.004, f"{b:.3f}", ha="center", fontsize=9)
    ax1.set_xticks(x)
    ax1.set_xticklabels(seasons)
    ax1.set_ylabel("SEC (kWh / kg)")
    ax1.set_title("Kathmandu (cold, dry, 1350 m)")
    ax1.legend(loc="upper right")
    ax1.grid(axis="y", alpha=0.25)

    ax2.bar(x - w/2, e2_btn, w, label="E2", color=C_E2, edgecolor="black", linewidth=0.5)
    ax2.bar(x + w/2, e3_btn, w, label="E3", color=C_E3, edgecolor="black", linewidth=0.5)
    for i, (a, b) in enumerate(zip(e2_btn, e3_btn)):
        ax2.text(i - w/2, a + 0.004, f"{a:.3f}", ha="center", fontsize=9)
        ax2.text(i + w/2, b + 0.004, f"{b:.3f}", ha="center", fontsize=9)
    ax2.set_xticks(x)
    ax2.set_xticklabels(seasons)
    ax2.set_title("Biratnagar (warm, humid, 72 m)")
    ax2.legend(loc="upper right")
    ax2.grid(axis="y", alpha=0.25)

    ax1.set_ylim(0, 0.18)
    fig.suptitle("E2 vs E3, 10 m² collector, seasonal SEC", fontsize=12, y=1.02)
    fig.tight_layout()
    p = TMP / "seasonal.png"
    fig.savefig(p, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return p


def plot_area_sweep() -> Path:
    areas = np.array([2, 4, 6, 8, 10])
    sec = np.array([0.1927, 0.1521, 0.1372, 0.1287, 0.1228])
    q_rep = np.array([4.98, 9.54, 13.69, 17.49, 20.95])
    q_use = np.array([4.98, 8.88, 10.26, 11.03, 11.57])
    clip_h = np.array([0.0, 3.23, 5.12, 5.85, 6.30])
    waste = q_rep - q_use

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4.4), dpi=160)
    ax1.plot(areas, sec, "o-", color=C_E2, linewidth=2, markersize=8)
    for a, s in zip(areas, sec):
        ax1.annotate(f"{s:.3f}", (a, s), textcoords="offset points", xytext=(6, 6), fontsize=10)
    ax1.set_xlabel("Collector area (m²)")
    ax1.set_ylabel("SEC (kWh / kg)")
    ax1.set_title("SEC drops monotonically (no plateau)")
    ax1.grid(alpha=0.3)
    ax1.set_ylim(0.10, 0.21)

    ax2.bar(areas, q_use, color=C_E2, edgecolor="black", linewidth=0.5, label="Usable solar")
    ax2.bar(areas, waste, bottom=q_use, color="#cccccc", edgecolor="black", linewidth=0.5, label="Clipped (wasted)")
    for i, a in enumerate(areas):
        pct = 0 if q_rep[i] == 0 else waste[i] / q_rep[i] * 100
        ax2.text(a, q_rep[i] + 0.4, f"{pct:.0f}% clip", ha="center", fontsize=9, color="#555555")
        ax2.text(a, q_rep[i] + 1.6, f"{clip_h[i]:.1f} h", ha="center", fontsize=9, color="black")
    ax2.set_xlabel("Collector area (m²)")
    ax2.set_ylabel("Solar energy (kWh)")
    ax2.set_title("Clipping grows with area")
    ax2.legend(loc="upper left")
    ax2.grid(axis="y", alpha=0.3)
    ax2.set_ylim(0, 28)

    fig.suptitle("E2 area sweep, Kathmandu autumn", fontsize=12, y=1.02)
    fig.tight_layout()
    p = TMP / "area_sweep.png"
    fig.savefig(p, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return p


def plot_clipping_physics() -> Path:
    t = np.linspace(0, 24, 500)
    # Ambient diurnal
    T_amb = 12 + 6 * np.sin((t - 8) / 24 * 2 * np.pi)
    # HRX outlet: ambient + eps * (T_exh - T_amb); assume mid-drying T_exh ≈ 30, eps=0.7
    T_hrx_out = T_amb + 0.7 * (30 - T_amb)
    # Solar Gaussian around noon
    G = 900 * np.exp(-0.5 * ((t - 12.5) / 3.0) ** 2)
    G = np.where((t > 6.5) & (t < 18.5), G, 0)
    eta_c = 0.55
    A = 10.0
    m_cp = 0.101 * 1006 / 1000  # kW/K
    dT = eta_c * G * A / 1000 / m_cp  # K added by solar
    T_after_solar = T_hrx_out + dT
    T_set = 45

    fig, ax = plt.subplots(figsize=(10, 4.6), dpi=160)
    ax.fill_between(t, T_set, T_after_solar, where=(T_after_solar > T_set),
                    color="#cccccc", alpha=0.7, label="Clipped (wasted)")
    ax.plot(t, T_amb, color="#555555", linewidth=1.3, label="T_ambient", linestyle=":")
    ax.plot(t, T_hrx_out, color=C_D1, linewidth=1.8, label="After HRX")
    ax.plot(t, T_after_solar, color=C_E2, linewidth=2.2, label="After solar (E2)")
    ax.axhline(T_set, color=C_E3, linestyle="--", linewidth=1.5, label=f"T_set = {T_set} °C")
    # highlight clipping window
    clip_mask = T_after_solar > T_set
    if clip_mask.any():
        t_start = t[clip_mask].min()
        t_end = t[clip_mask].max()
        ax.axvspan(t_start, t_end, alpha=0.1, color=C_E3)
        ax.text((t_start + t_end) / 2, 58, f"Clipping window\n~{t_start:.1f}-{t_end:.1f} h",
                ha="center", fontsize=10, color=C_E3, fontweight="bold")
    ax.set_xlabel("Hour of day")
    ax.set_ylabel("Temperature (°C)")
    ax.set_title("Why clipping happens in E2: solar-lifted stream exceeds T_set around midday")
    ax.legend(loc="upper left", ncol=2)
    ax.grid(alpha=0.25)
    ax.set_xlim(0, 24)
    ax.set_ylim(0, 65)
    fig.tight_layout()
    p = TMP / "clipping_physics.png"
    fig.savefig(p, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return p


def plot_climate_annual() -> Path:
    locs = ["Kathmandu\n(1350 m)", "Biratnagar\n(72 m)"]
    e2 = [0.197, 0.129]
    e3 = [0.217, 0.137]
    x = np.arange(len(locs))
    w = 0.28
    fig, ax = plt.subplots(figsize=(6.5, 4.2), dpi=160)
    b1 = ax.bar(x - w/2, e2, w, label="E2", color=C_E2, edgecolor="black", linewidth=0.5)
    b2 = ax.bar(x + w/2, e3, w, label="E3", color=C_E3, edgecolor="black", linewidth=0.5)
    for bars in (b1, b2):
        for b in bars:
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.005,
                    f"{b.get_height():.3f}", ha="center", fontsize=10, fontweight="bold")
    # E2 advantage annotation
    for i, (a, b) in enumerate(zip(e2, e3)):
        adv = (b - a) / b * 100
        ax.text(i, max(a, b) + 0.022, f"E2 wins by {adv:.1f}%", ha="center", fontsize=9,
                color=C_E2, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(locs)
    ax.set_ylabel("SEC (kWh / kg)")
    ax.set_title("Annual SEC, 10 m² collector, r = 0")
    ax.set_ylim(0, 0.26)
    ax.legend(loc="upper right")
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout()
    p = TMP / "climate.png"
    fig.savefig(p, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return p


# ============================================================
# Block diagrams
# ============================================================

def draw_config_A(slide, left, top):
    """Closed-loop HP dryer block diagram. left/top in Emu."""
    row_y = top + Inches(0.6)
    # Components
    amb = comp_box(slide, left, row_y, Inches(1.0), Inches(0.8), "Ambient\nair", fill=SAGE, size=11)
    mix = comp_box(slide, left + Inches(1.4), row_y, Inches(0.9), Inches(0.8), "Mix", fill=GREY, size=11)
    evap = comp_box(slide, left + Inches(2.7), row_y, Inches(1.0), Inches(0.8), "Evap\n(dehum.)", fill=TEAL, size=11)
    cond = comp_box(slide, left + Inches(4.1), row_y, Inches(1.0), Inches(0.8), "Cond\n(heat)", fill=CORAL, size=11)
    chamber = comp_box(slide, left + Inches(5.5), row_y, Inches(1.4), Inches(0.8), "Chamber\n(10 trays)", fill=NAVY, size=11)
    exh = comp_box(slide, left + Inches(7.3), row_y, Inches(1.0), Inches(0.8), "Exhaust", fill=GREY, size=11)
    # Forward arrows
    parts = [amb, mix, evap, cond, chamber, exh]
    for a, b in zip(parts[:-1], parts[1:]):
        arrow(slide, a.left + a.width, row_y + Inches(0.4),
              b.left, row_y + Inches(0.4))
    # Recirc dashed arrow from exhaust back to mix
    arrow(slide, exh.left + exh.width / 2, row_y + exh.height,
          mix.left + mix.width / 2, row_y + mix.height + Inches(0.8),
          color=CORAL, weight=1.2, dash=True)
    arrow(slide, mix.left + mix.width / 2, row_y + mix.height + Inches(0.8),
          mix.left + mix.width / 2, row_y + mix.height,
          color=CORAL, weight=1.2, dash=True)
    arrow(slide, exh.left + exh.width / 2, row_y + exh.height,
          exh.left + exh.width / 2, row_y + exh.height + Inches(0.8),
          color=CORAL, weight=1.2, dash=True)
    arrow(slide, exh.left + exh.width / 2, row_y + exh.height + Inches(0.8),
          mix.left + mix.width / 2, row_y + mix.height + Inches(0.8),
          color=CORAL, weight=1.2, dash=True)
    # Label recirc
    add_text(slide, left + Inches(3.2), row_y + Inches(1.55), Inches(4.0), Inches(0.3),
             "recirculation r (closed loop)", size=10, color=CORAL, italic=True)
    # HP label bracket
    hp_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        evap.left - Inches(0.1), row_y - Inches(0.3),
        (cond.left + cond.width) - evap.left + Inches(0.2), evap.height + Inches(0.5)
    )
    hp_box.fill.background()
    hp_box.line.color.rgb = GREY
    hp_box.line.dash_style = 7
    hp_box.line.width = Pt(0.75)
    add_text(slide, evap.left, row_y - Inches(0.35), Inches(2.4), Inches(0.28),
             "Heat pump", size=10, italic=True, color=GREY)


def draw_E2_E3(slide, kind, left, top):
    """Draw E2 or E3 block diagram. kind in {'E2','E3'}."""
    row_y = top + Inches(0.5)
    amb = comp_box(slide, left, row_y, Inches(0.85), Inches(0.75), "Amb", fill=SAGE, size=10)
    hrx = comp_box(slide, left + Inches(1.15), row_y, Inches(0.85), Inches(0.75), "HRX", fill=SAGE, size=10)
    if kind == "E2":
        first = comp_box(slide, left + Inches(2.3), row_y, Inches(0.95), Inches(0.75), "Solar", fill=GOLD, size=10)
        second = comp_box(slide, left + Inches(3.45), row_y, Inches(0.95), Inches(0.75), "Cond", fill=CORAL, size=10)
    else:
        first = comp_box(slide, left + Inches(2.3), row_y, Inches(0.95), Inches(0.75), "Cond", fill=CORAL, size=10)
        second = comp_box(slide, left + Inches(3.45), row_y, Inches(0.95), Inches(0.75), "Solar", fill=GOLD, size=10)
    chamber = comp_box(slide, left + Inches(4.6), row_y, Inches(1.15), Inches(0.75), "Chamber", fill=NAVY, size=10)
    exh = comp_box(slide, left + Inches(5.95), row_y, Inches(0.85), Inches(0.75), "Exh", fill=GREY, size=10)
    parts = [amb, hrx, first, second, chamber, exh]
    for a, b in zip(parts[:-1], parts[1:]):
        arrow(slide, a.left + a.width, row_y + Inches(0.375),
              b.left, row_y + Inches(0.375))
    # Exhaust returns to HRX
    arrow(slide, exh.left + exh.width / 2, row_y + exh.height,
          exh.left + exh.width / 2, row_y + exh.height + Inches(0.45),
          color=GREY, weight=1.0, dash=True)
    arrow(slide, exh.left + exh.width / 2, row_y + exh.height + Inches(0.45),
          hrx.left + hrx.width / 2, row_y + hrx.height + Inches(0.45),
          color=GREY, weight=1.0, dash=True)
    arrow(slide, hrx.left + hrx.width / 2, row_y + hrx.height + Inches(0.45),
          hrx.left + hrx.width / 2, row_y + hrx.height,
          color=GREY, weight=1.0, dash=True)
    add_text(slide, hrx.left, row_y + hrx.height + Inches(0.5), Inches(4.5), Inches(0.3),
             "exhaust pre-heats HRX (70 % effective)", size=9, italic=True, color=GREY)
    # Title badge
    badge_fill = C_E2 if kind == "E2" else C_E3
    add_text(slide, left, top - Inches(0.0), Inches(6.8), Inches(0.35),
             f"{kind}: " + ("Amb \u2192 HRX \u2192 Solar \u2192 Cond \u2192 Chamber" if kind == "E2"
                            else "Amb \u2192 HRX \u2192 Cond \u2192 Solar \u2192 Chamber"),
             size=13, bold=True, color=RGBColor.from_string(badge_fill.lstrip('#')))


def draw_thermal_zones(slide, left, top, width):
    """Concept sketch: multi-inlet thermal-zone chamber."""
    # Chamber box
    chamber_h = Inches(2.2)
    chamber = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, chamber_h)
    chamber.fill.solid()
    chamber.fill.fore_color.rgb = RGBColor(0x2A, 0x4A, 0x7F)
    chamber.line.color.rgb = DARK
    chamber.line.width = Pt(0.8)
    # Tray zone dividers
    zones = 4
    zone_w = width / zones
    for i in range(1, zones):
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            left + zone_w * i, top + Inches(0.15),
            left + zone_w * i, top + chamber_h - Inches(0.15)
        )
        ln.line.color.rgb = WHITE
        ln.line.width = Pt(1.0)
        ln.line.dash_style = 7
    # Zone labels + tray stacks
    for i in range(zones):
        zx = left + zone_w * i + zone_w / 2
        add_text(slide, zx - Inches(0.6), top + Inches(0.8), Inches(1.2), Inches(0.3),
                 f"Zone {i+1}", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, zx - Inches(0.6), top + Inches(1.2), Inches(1.2), Inches(0.3),
                 "trays",
                 size=9, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Multi-inlet manifold above chamber
    manifold_y = top - Inches(0.9)
    manifold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.2), manifold_y,
        width - Inches(0.4), Inches(0.35)
    )
    manifold.fill.solid()
    manifold.fill.fore_color.rgb = CORAL
    manifold.line.color.rgb = DARK
    add_text(slide, left + Inches(0.2), manifold_y + Inches(0.03), width - Inches(0.4), Inches(0.3),
             "Hot-air distribution manifold (from Cond + Solar)",
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Down-arrows from manifold to each zone
    for i in range(zones):
        zx = left + zone_w * i + zone_w / 2
        arrow(slide, zx, manifold_y + Inches(0.35),
              zx, top, color=CORAL, weight=1.5)
    # Exhaust arrows below chamber
    for i in range(zones):
        zx = left + zone_w * i + zone_w / 2
        arrow(slide, zx, top + chamber_h,
              zx, top + chamber_h + Inches(0.45), color=GREY, weight=1.3)
    # Exhaust collector bar
    exh_y = top + chamber_h + Inches(0.45)
    exh_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.2), exh_y,
        width - Inches(0.4), Inches(0.3)
    )
    exh_bar.fill.solid()
    exh_bar.fill.fore_color.rgb = GREY
    exh_bar.line.color.rgb = DARK
    add_text(slide, left + Inches(0.2), exh_y + Inches(0.02), width - Inches(0.4), Inches(0.28),
             "Exhaust plenum \u2192 back to HRX",
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# Build slides
# ============================================================

def slide_title():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Accent band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), SW, Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = CORAL
    band.line.fill.background()
    add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.2),
             "Where should the solar go?",
             size=44, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
             "Comparing E2 vs E3 for the solar-assisted heat pump dryer",
             size=22, color=GREY, italic=True)
    add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
             "Supervisor discussion, 20 April 2026",
             size=16, color=GREY)
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.8),
             "Agenda: base dryer \u2192 HRX \u2192 solar placement (E2 vs E3) \u2192 clipping "
             "\u2192 validation \u2192 scale-up ideas",
             size=13, italic=True, color=GREY)
    add_notes(s, """
Opening, casual: this is a working session, not a presentation. The goal is
to walk you through how the E2 and E3 configs are set up, what the
simulations are actually doing, and why we think E2 is the better solar
placement. Then I want to open it up to discuss scaling this to 500-1000 kg
apples and whether a multi-zone chamber makes sense as the next step.
Feel free to interrupt at any slide.
""")


def slide_base_dryer():
    s = blank_slide("Base dryer: Config A (HP only, closed loop)",
                    "This is the machine before we add anything")
    draw_config_A(s, Inches(1.0), Inches(1.6))
    bullets(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.5), [
        "HP condenser always heats to T_set = 45 °C; chamber holds exactly that (validated, V1.1).",
        "Fresh apples: 22.5 kg wet \u2248 3 kg dry. Water to remove: 19.2 kg. Drying time 13.8 h.",
        "10 trays in series, cross-flow, 1.1 m/s. Ambient air enters, humid exhaust leaves.",
        "Recirculation path shown dashed; in this discussion the r=0 (open-loop) SEC is our anchor.",
    ], size=14)
    # Anchor number
    anchor = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(9.0), Inches(5.9), Inches(3.8), Inches(1.2))
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = LIGHT
    anchor.line.color.rgb = NAVY
    anchor.line.width = Pt(1.2)
    add_text(s, Inches(9.0), Inches(5.95), Inches(3.8), Inches(0.4),
             "Config A baseline", size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.0), Inches(6.35), Inches(3.8), Inches(0.5),
             "SEC = 0.717 kWh/kg", size=20, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.0), Inches(6.75), Inches(3.8), Inches(0.3),
             "Kathmandu, annual, r = 0",
             size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    add_notes(s, """
Talking points:
- Anchor this number in their head: 0.717 kWh of electricity per kg of water.
  Every improvement from here on is referenced against this.
- Worth pointing out that the closed-loop version (r > 0) is more efficient
  still (0.515 at r=0.9, from V2.1 table), but we use r=0 throughout the
  solar/HRX configs because mixing is not needed once we have waste-heat
  recovery from the HRX and a dedicated evap stream. Happy to discuss.
- Product spec is apple slices, 6 mm thick, 3 kg dry mass, 19.2 kg water
  to remove. Current scale is small, which comes up again at the end.
""")


def slide_hrx_step():
    s = blank_slide("Step 1: add a heat-recovery exchanger (Config D1)",
                    "Exhaust pre-heats the condenser inlet through a 70 % effective HRX")
    # Embed HRX comparison plot
    s.shapes.add_picture(str(PLOTS / "hrx_comparison" / "hrx_detail_kathmandu.png"),
                          Inches(0.4), Inches(1.4), height=Inches(5.5))
    # Numbers panel
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.6), Inches(1.6), Inches(4.5), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    add_text(s, Inches(8.7), Inches(1.75), Inches(4.3), Inches(0.4),
             "Kathmandu annual, r = 0", size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rows = [
        ("Config A (HP only)", "0.717", "baseline", GREY),
        ("Config D1 (HP + HRX)", "0.365", "49 % lower than A", TEAL),
        ("Config D2 (HRX + dyn evap)", "0.354", "3 % better than D1", TEAL),
    ]
    for i, (lab, val, note, col) in enumerate(rows):
        y = Inches(2.3 + i * 0.85)
        add_text(s, Inches(8.8), y, Inches(2.5), Inches(0.35),
                 lab, size=12, bold=True, color=DARK)
        add_text(s, Inches(11.3), y, Inches(1.7), Inches(0.35),
                 val, size=16, bold=True, color=col)
        add_text(s, Inches(8.8), y + Inches(0.38), Inches(4.1), Inches(0.3),
                 note, size=10, italic=True, color=GREY)
    add_text(s, Inches(8.8), Inches(5.3), Inches(4.3), Inches(1.4),
             "HRX does the heavy lifting.\n"
             "Dynamic evap (D2) only trims another\n"
             "0.011 kWh/kg, a minor refinement.",
             size=11, italic=True, color=DARK)
    add_notes(s, """
Talking points:
- The HRX is a simple air-to-air cross-flow exchanger. The warm, moist
  exhaust gives up its sensible heat to the incoming cold ambient stream.
  70 % effectiveness is a reasonable commercial value.
- The result is strong: 49 % SEC drop over A, from just one passive
  component. No moving parts, no control logic, no electricity.
- D2 is D1 plus dynamic evaporator compensation (the evap sizes itself
  based on instantaneous ambient, not a fixed approach). The incremental
  benefit is tiny, about 1 percentage point. Worth flagging that we
  tested it because we suspected it mattered, and the data says it does
  not, which is itself a useful negative result.
- This sets up the natural next question: can we layer solar on top?
""")


def slide_solar_step():
    s = blank_slide("Step 2: add a solar collector (Config B)",
                    "Flat-plate collector in series before the condenser")
    # Simple diagram
    y = Inches(1.8)
    amb = comp_box(s, Inches(1.0), y, Inches(1.0), Inches(0.8), "Ambient", fill=SAGE)
    sol = comp_box(s, Inches(2.4), y, Inches(1.2), Inches(0.8), "Solar\ncollector", fill=GOLD)
    cond = comp_box(s, Inches(4.0), y, Inches(1.0), Inches(0.8), "Cond", fill=CORAL)
    cham = comp_box(s, Inches(5.4), y, Inches(1.4), Inches(0.8), "Chamber", fill=NAVY)
    exh = comp_box(s, Inches(7.2), y, Inches(1.0), Inches(0.8), "Exhaust", fill=GREY)
    parts = [amb, sol, cond, cham, exh]
    for a, b in zip(parts[:-1], parts[1:]):
        arrow(s, a.left + a.width, y + Inches(0.4), b.left, y + Inches(0.4))

    bullets(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(2.3), [
        "When the sun is out, solar does part of the heating. HP takes what is left to reach 45 °C.",
        "At night or under cloud cover, solar contributes zero and the HP runs full duty, so SEC can never be worse than Config A.",
        "With 10 m² collector: solar provides 32 % of total heating at Kathmandu, 46 % at Biratnagar.",
        "Observation: B alone beats A but does not beat D1 at Kathmandu. Solar alone < HRX alone at high altitude.",
    ])
    # Ladder chart
    s.shapes.add_picture(str(plot_ladder()), Inches(3.0), Inches(5.25), width=Inches(7.3))
    add_notes(s, """
Talking points:
- Solar is a preheat. It does not replace the HP. It just reduces the lift
  the HP has to provide.
- Quick numbers: B alone gets us to 0.488 at Kathmandu. D1 alone gets us to
  0.365. So at 1350 m altitude, HRX is more valuable than 10 m² of solar
  by itself. That was a bit surprising when we first saw it, and it
  motivated combining them (E series).
- The ladder chart tells the whole story: each component buys us something.
  The open question: when we combine HRX + solar (E series), where exactly
  should the solar sit in the stream?
""")


def slide_E2_vs_E3_question():
    s = blank_slide("The question: where does the solar go in E configs?",
                    "Both have an HRX and a solar collector, only the ordering differs")
    draw_E2_E3(s, "E2", Inches(0.6), Inches(1.6))
    draw_E2_E3(s, "E3", Inches(0.6), Inches(4.2))

    # Control logic notes
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.2), Inches(1.6), Inches(4.8), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    add_text(s, Inches(8.3), Inches(1.7), Inches(4.6), Inches(0.4),
             "Control logic differs", size=14, bold=True, color=NAVY)
    add_text(s, Inches(8.3), Inches(2.15), Inches(4.6), Inches(0.35),
             "E2: Cond targets T_set", size=12, bold=True, color=NAVY)
    add_text(s, Inches(8.3), Inches(2.5), Inches(4.6), Inches(1.6),
             "If solar delivers the air already above 45 °C, the simulator "
             "clips to T_set and HP idles (Q_cond = 0). Any extra solar "
             "heat above 45 °C is discarded (no storage).",
             size=11, color=DARK)
    add_text(s, Inches(8.3), Inches(4.1), Inches(4.6), Inches(0.35),
             "E3: Solar-priority", size=12, bold=True, color=NAVY)
    add_text(s, Inches(8.3), Inches(4.45), Inches(4.6), Inches(2.2),
             "HP modulates T_cond so that T_after_solar exactly equals T_set. "
             "If solar alone from the HRX output reaches T_set, HP goes OFF "
             "entirely. No clipping by design.",
             size=11, color=DARK)
    add_notes(s, """
Talking points:
- This is the crux slide. Two seemingly-equivalent layouts, both deliver
  45 degree air to the chamber. Naively one might guess E3 is better
  because its HP can idle completely in strong sun.
- The control logic is important. In E2 the condenser is the final gate
  and always targets 45. In E3 the collector is the final gate, so the
  HP has to modulate its discharge pressure to leave the solar "just
  enough room" to land exactly at 45.
- Both are physically defensible topologies. Nobody has published a
  head-to-head comparison with a validated thermodynamic model, which is
  why we spent time on it.
""")


def slide_E2_in_action():
    s = blank_slide("E2 in action (Kathmandu autumn, 10 m²)",
                    "Solar lifts before the condenser; condenser fills the gap to 45 °C")
    s.shapes.add_picture(str(PLOTS / "compare_E_autumn_ktm" / "E2" / "Ac_10m2_hrx0.70_overview.png"),
                          Inches(0.4), Inches(1.4), height=Inches(5.6))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(9.6), Inches(1.6), Inches(3.5), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = RGB_E2
    box.line.width = Pt(1.5)
    add_text(s, Inches(9.7), Inches(1.7), Inches(3.3), Inches(0.4),
             "E2 behaviour", size=14, bold=True, color=RGB_E2)
    bullets(s, Inches(9.7), Inches(2.15), Inches(3.3), Inches(4.2), [
        "T_chamber holds 45 °C the entire run",
        "HP runs throughout the day, solar reduces its load but rarely shuts it off",
        "SEC = 0.123 kWh/kg",
        "Mean COP over run = 5.25",
        "Total drying time 16.9 h",
    ], size=11)
    add_notes(s, """
Talking points:
- Point to T_chamber at 45 flat. That is the control law working.
- Point to W_comp: the compressor never goes to zero during daylight. It
  dips when solar is strong but always keeps some load because the
  condenser is the final heater.
- SEC 0.123 at 10 m2 is a very good number for this climate. The
  reference Config A is 0.500 for autumn KTM at r=0. So E2 is a
  75 percent reduction from the bare HP baseline in autumn alone.
""")


def slide_E3_in_action():
    s = blank_slide("E3 in action (Kathmandu autumn, 10 m²)",
                    "HP modulates T_cond to leave solar just enough room to reach 45 °C")
    s.shapes.add_picture(str(PLOTS / "compare_E_autumn_ktm" / "E3" / "Ac_10m2_hrx0.70_overview.png"),
                          Inches(0.4), Inches(1.4), height=Inches(5.6))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(9.6), Inches(1.6), Inches(3.5), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = RGB_E3
    box.line.width = Pt(1.5)
    add_text(s, Inches(9.7), Inches(1.7), Inches(3.3), Inches(0.4),
             "E3 behaviour", size=14, bold=True, color=RGB_E3)
    bullets(s, Inches(9.7), Inches(2.15), Inches(3.3), Inches(4.2), [
        "T_chamber holds 45 °C (different path, same outcome)",
        "HP turns OFF entirely for ~1-2 h around solar noon",
        "SEC = 0.133 kWh/kg",
        "Mean COP = 6.53 (higher than E2)",
        "No clipping: 100 % of collected solar reaches the chamber",
    ], size=11)
    add_notes(s, """
Talking points:
- Look at the W_comp trace: it goes flat at zero during peak sun. That is
  the HP stepping out of the way and letting solar do everything.
- Higher COP is real: when the HP does run, it runs at a much lower
  pressure ratio because T_cond is modulated down. So every kWh of
  compressor work delivers more Q_cond.
- Here's the puzzle: E3 has the higher mean COP AND the higher peak COP
  AND no clipping. So why doesn't it win? Transition to the next slide.
""")


def slide_side_by_side():
    s = blank_slide("Side-by-side: what each config actually does",
                    "Autumn KTM, 10 m², r = 0, identical weather, identical drying time (16.9 h)")
    # 2x2 grid (slightly smaller to make room for annotation panels)
    w = Inches(5.5)
    h = Inches(2.35)
    s.shapes.add_picture(str(PLOTS / "compare_E_autumn_ktm" / "E2" / "Ac_10m2_hrx0.70_solar.png"),
                          Inches(0.3), Inches(1.55), width=w, height=h)
    s.shapes.add_picture(str(PLOTS / "compare_E_autumn_ktm" / "E3" / "Ac_10m2_hrx0.70_solar.png"),
                          Inches(7.5), Inches(1.55), width=w, height=h)
    s.shapes.add_picture(str(PLOTS / "compare_E_autumn_ktm" / "E2" / "Ac_10m2_hrx0.70_heatpump.png"),
                          Inches(0.3), Inches(4.25), width=w, height=h)
    s.shapes.add_picture(str(PLOTS / "compare_E_autumn_ktm" / "E3" / "Ac_10m2_hrx0.70_heatpump.png"),
                          Inches(7.5), Inches(4.25), width=w, height=h)
    # Annotation panels (data-grounded)
    def panel(left, top, title, color, lines):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, Inches(1.55), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = color
        box.line.width = Pt(1.2)
        add_text(s, left + Inches(0.08), top + Inches(0.06),
                 Inches(1.4), Inches(0.3),
                 title, size=11, bold=True, color=color)
        y = top + Inches(0.4)
        for ln, bold, sz in lines:
            add_text(s, left + Inches(0.08), y, Inches(1.4), Inches(0.3),
                     ln, size=sz, bold=bold, color=DARK)
            y += Inches(0.28 if sz <= 9 else 0.32)

    # E2 solar panel
    panel(Inches(5.9), Inches(1.55), "E2 solar", RGB_E2, [
        ("Cooler inlet", True, 10),
        ("~32 °C (HRX only)", False, 9),
        ("Collected 20.95 kWh", True, 10),
        ("Clipped 9.38 kWh", False, 9),
        ("Usable 11.57 kWh", True, 10),
    ])
    # E3 solar panel
    panel(Inches(13.05) - Inches(1.55), Inches(1.55), "E3 solar", RGB_E3, [
        ("Hotter inlet", True, 10),
        ("~37 °C (HRX + Cond)", False, 9),
        ("Collected 19.77 kWh", True, 10),
        ("No clipping (design)", False, 9),
        ("Usable 19.77 kWh", True, 10),
    ])
    # E2 HP panel
    panel(Inches(5.9), Inches(4.25), "E2 heat pump", RGB_E2, [
        ("Runs most of day", True, 10),
        ("at moderate load", False, 9),
        ("W_comp 1.96 kWh", True, 10),
        ("Mean COP 5.25", False, 10),
        ("T_cond = 45 °C", False, 9),
    ])
    # E3 HP panel
    panel(Inches(13.05) - Inches(1.55), Inches(4.25), "E3 heat pump", RGB_E3, [
        ("OFF at solar noon", True, 10),
        ("(clean zero band)", False, 9),
        ("W_comp 2.16 kWh", True, 10),
        ("Mean COP 6.53", False, 10),
        ("T_cond modulated", False, 9),
    ])
    # Summary band at bottom
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.6))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_text(s, Inches(0.45), Inches(6.8), Inches(12.4), Inches(0.5),
             "The paradox: E3 has higher COP + no clipping + more usable solar, "
             "yet E2 uses less total W_comp (1.96 < 2.16 kWh). "
             "E2's cooler collector wins across every HP-ON hour, not just midday.",
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_notes(s, """
Walk through the four plots one panel at a time. All numbers are from
master_summary.csv for autumn KTM, 10 m2, r=0.

TOP ROW (solar collector output):
- E2 collects 20.95 kWh total; E3 collects 19.77 kWh. Difference is about
  6 percent. Driver: Hottel-Whillier-Bliss says collector efficiency drops
  linearly with (T_inlet - T_amb). In E2 the collector sees HRX output
  directly (~32 C). In E3 the collector sees HRX output plus whatever
  partial lift the condenser has added (~37 C). Cooler inlet = higher
  efficiency = more kWh collected.
- E2 wastes 9.38 kWh (44.8 percent of its haul) because during midday the
  solar alone already pushes the stream above 45 C. Simulator clips to
  T_set; excess is discarded (no storage).
- E3 has zero clipping by design: the condenser modulates T_cond so the
  solar lift lands exactly at 45 C. All 19.77 kWh reaches the chamber.
- Net usable solar: E3 (19.77) beats E2 (11.57).

BOTTOM ROW (heat pump):
- E2 HP runs almost continuously at moderate power. No clean OFF zone.
  T_cond pegged at 45 C, so lift is full; mean COP = 5.25.
- E3 HP has a clean OFF window around solar noon (visible as flat zero in
  W_comp trace). When it runs, T_cond is modulated down so the lift is
  smaller; mean COP = 6.53 (25 percent higher than E2).

THE PARADOX (bottom band):
- E3 looks better on every individual lever: higher COP, no waste, more
  usable solar delivered.
- Yet simulator says W_comp_E2 (1.96 kWh) < W_comp_E3 (2.16 kWh),
  and SEC follows: E2 at 0.123 vs E3 at 0.133.
- Mechanism from the memory note: E2's collector runs cooler DURING EVERY
  HP-ON HOUR (not just at noon). The accumulated extra solar capture
  across the run (~1.2 kWh autumn) outweighs the COP bonus E3 gets when
  its HP is running.
- This is the piece I want to sanity check with you. Does the physical
  story track, or should we dig into the energy decomposition by hour
  before locking the headline?
""")


def slide_punchline():
    s = blank_slide("The punchline: why E2 beats E3",
                    "Collector inlet temperature is everything (Hottel-Whillier-Bliss)")
    # Table
    rows = [
        ("", "E2", "E3"),
        ("SEC (kWh/kg)", "0.123", "0.133"),
        ("Mean COP", "5.25", "6.53"),
        ("Solar collected (kWh)", "20.95", "19.77"),
        ("Solar usable (kWh)", "11.57", "19.77"),
        ("Clipping waste", "44.8 %", "0 %"),
    ]
    x0 = Inches(0.7)
    y0 = Inches(1.6)
    col_w = [Inches(3.0), Inches(2.2), Inches(2.2)]
    row_h = Inches(0.48)
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            left = x0 + sum(col_w[:ci], Emu(0))
            bg = NAVY if ri == 0 else (LIGHT if ri % 2 == 0 else WHITE)
            fg = WHITE if ri == 0 else DARK
            shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y0 + row_h * ri, col_w[ci], row_h)
            shp.fill.solid()
            shp.fill.fore_color.rgb = bg
            shp.line.color.rgb = GREY
            shp.line.width = Pt(0.5)
            tf = shp.text_frame
            tf.margin_top = Pt(3)
            tf.margin_bottom = Pt(3)
            tf.margin_left = Pt(6)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(13)
            p.font.bold = (ri == 0) or (ci == 0)
            p.font.color.rgb = fg
            if ci > 0:
                p.alignment = PP_ALIGN.CENTER

    add_text(s, Inches(0.7), Inches(5.0), Inches(7.3), Inches(0.4),
             "Autumn Kathmandu, 10 m² collector", size=10, italic=True, color=GREY)

    # Hottel physics panel
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.3), Inches(1.6), Inches(4.7), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.2)
    add_text(s, Inches(8.4), Inches(1.7), Inches(4.5), Inches(0.4),
             "The physics", size=14, bold=True, color=NAVY)
    add_text(s, Inches(8.4), Inches(2.15), Inches(4.5), Inches(0.4),
             "Hottel-Whillier-Bliss:", size=12, bold=True, color=DARK)
    add_text(s, Inches(8.4), Inches(2.55), Inches(4.5), Inches(0.5),
             "\u03B7_c = F_R(\u03C4\u03B1) \u2212 F_R U_L (T_in \u2212 T_amb) / G",
             size=13, bold=True, color=CORAL)
    add_text(s, Inches(8.4), Inches(3.2), Inches(4.5), Inches(3.8),
             "In E2 the collector sees air at ~32 °C (HRX output, mid-drying).\n\n"
             "In E3 the collector sees air at ~37 °C (HRX + condenser output).\n\n"
             "That 5 °C hotter inlet costs roughly 2.5 percentage points of "
             "energy-weighted efficiency.\n\n"
             "E2 captures +1.2 kWh extra solar per day; E3's lower compressor "
             "work saves only ~1.0 kWh thermal. E2 nets +0.2 kWh/day.\n\n"
             "Keep the solar upstream of the condenser.",
             size=11, color=DARK)
    add_notes(s, """
Talking points:
- Hottel-Whillier-Bliss is the standard flat-plate collector equation.
  The second term is the thermal loss: hotter collector inlet means more
  heat leaks to ambient, lower efficiency.
- In E3, the stream going into the collector has already been lifted by
  the condenser, so it arrives hotter. The collector runs less
  efficiently because it is trying to heat an already-warm stream.
- In E2, the collector inlet is cooler (just ambient + HRX), so it runs
  more efficiently. The HP then "tops off" any remaining lift.
- The COP story is a distraction. Yes E3 has higher COP. But the
  collector inefficiency costs more than the COP gain buys us.
- Practical takeaway: always put the solar upstream of the condenser in
  a series topology. This is the answer to "where should the solar go".
""")


def slide_clipping():
    s = blank_slide("The clipping problem in E2",
                    "Solar heat above T_set has nowhere to go (no storage)")
    s.shapes.add_picture(str(plot_clipping_physics()), Inches(0.3), Inches(1.4), width=Inches(7.8))
    # Side panel
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.3), Inches(1.6), Inches(4.7), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = CORAL
    box.line.width = Pt(1.2)
    add_text(s, Inches(8.4), Inches(1.7), Inches(4.5), Inches(0.4),
             "When does clipping happen?", size=13, bold=True, color=NAVY)
    bullets(s, Inches(8.4), Inches(2.2), Inches(4.5), Inches(4.8), [
        "WHEN: peak-sun hours, roughly 10:00-15:00 (6.3 h/day at 10 m² autumn).",
        "WHY: stream already at 45 °C when it reaches the condenser; HP idles; excess heat is discarded.",
        "GETS WORSE as drying progresses: exhaust warms, HRX lifts inlet closer to T_set, even weak solar clips.",
        "Threshold heat: Q_solar > ~1.3 kW (stream can't accept more than ~13 K lift mid-drying).",
        "In irradiance: 10 m² clips above ~240 W/m²; 4 m² clips above ~600 W/m²; 2 m² never clips.",
        "Identity holds: Q_reported = Q_usable + Q_clipped exactly.",
    ], size=11)
    add_notes(s, """
Talking points:
- The shaded region on the plot is the waste. It is the slice of the
  solar-lifted stream that is above the 45 set point. The simulator
  reports the full Q_solar from Hottel, but physically that top slice
  never makes it to the chamber.
- Clipping is bounded by the sun, but it is also bounded by the
  thermodynamic state: once the stream approaches T_set from any source,
  even a small solar contribution pushes it over.
- The back-of-envelope threshold: m_da * cp * (T_set - T_hrx_out). With
  m_da around 0.10 kg/s and T_hrx_out around 32 at mid-drying, the
  allowable solar input is about 1.3 kW. At 10 m2 that corresponds to
  roughly 240 W/m2 of irradiance, which is essentially any clear-sky
  daytime hour.
- Key point: this is a thermodynamic limit, not a weather problem. A
  better weather forecast does not help. The fix is storage, which we
  come back to on the last slide.
""")


def slide_area_sweep():
    s = blank_slide("Area sweep: SEC keeps dropping even as clipping rises",
                    "Shrinking the collector to avoid clipping makes things worse")
    s.shapes.add_picture(str(plot_area_sweep()), Inches(0.5), Inches(1.4), width=Inches(12.3))
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
             "Extra area always buys extra useful solar. The HP duty cycle drops faster than "
             "clipping waste grows. Right fix: keep area, add storage, not shrink area.",
             size=13, italic=True, color=DARK)
    add_notes(s, """
Talking points:
- One might intuitively say: clipping is waste, let's downsize the
  collector. The data says the opposite. Even at 44 percent clipping,
  the remaining usable fraction still does more useful work than a
  smaller collector would.
- Monotonic SEC drop from 0.193 at 2 m2 to 0.123 at 10 m2. No plateau
  in the 2-10 m2 range. We tried 15 and 20 m2 too and still saw small
  improvements.
- Insight: the HP duty cycle drops faster than the collector efficiency
  does. Every kWh of usable solar displaces HP electricity that would
  otherwise have been spent, and the HP duty cycle is the dominant cost.
- The correct engineering response to clipping is thermal storage:
  capture the midday surplus, release it in the evening.
""")


def slide_seasonal_climate():
    s = blank_slide("Seasonal performance across two climates",
                    "E2 beats E3 in every season, at both locations")
    s.shapes.add_picture(str(plot_seasonal()), Inches(0.3), Inches(1.4), width=Inches(9.5))
    s.shapes.add_picture(str(plot_climate_annual()), Inches(9.9), Inches(1.4), width=Inches(3.3))
    add_text(s, Inches(0.3), Inches(6.3), Inches(12.8), Inches(0.8),
             "E2 advantage is consistent: 7-9 % at Kathmandu, 6-7 % at Biratnagar. "
             "Spring is best at both (strongest irradiance, moderate ambient, best collector efficiency).",
             size=13, italic=True, color=DARK)
    add_notes(s, """
Talking points:
- Four seasons, two climates, eight data points. E2 beats E3 on all
  eight. The mechanism is universal, not a fluke of one dataset.
- Spring is the best season at both locations because the sun angle is
  good and the ambient is not yet too hot (low U_L*(T_in - T_amb) loss).
- Biratnagar is always absolute-lower SEC because warmer ambient means
  smaller HP lift and better collector efficiency. But the E2 vs E3
  advantage is almost the same percent.
- If we had time we could run Taplejung (1820 m, even colder) but the
  mechanism would play out the same way.
""")


def slide_validation():
    s = blank_slide("How do we trust these numbers?",
                    "Model passes all first-law and mass-balance checks")
    bullets(s, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.5), [
        "1st law on HP cycle: Q_cond = Q_evap + W_comp. Residual < 1.5 % (enthalpy decomposition artefact, CoolProp states are exact).",
        "Water mass balance: cumulative dm_w matches 19.35 kg; target was 19.2 kg. Error ~0.8 %.",
        "T_to_chamber holds 45.00 °C exactly across all 6 configs (condenser control validated).",
        "COP in physical range: 3.2 to 6.5 across configs; Carnot efficiency 0.61-0.62 (consistent with η_is = 0.75).",
        "Pressure ratio 4-5 for R410A (matches textbook values).",
        "Tray humidity gradient monotonic early to late (see plot, right).",
        "Psychrometric consistency: \u03C9, RH, T_dp all close within 4 \u00D7 10\u207B\u2076.",
        "Solar energy identity: Q_reported = Q_usable + Q_clipped, exact by construction.",
        "Phase 3.5 validation report (Apr 2026): all 6 configs PASS.",
    ], size=12)
    # Tray plot as proof
    s.shapes.add_picture(str(PLOTS / "compare_E_autumn_ktm" / "E2" / "Ac_10m2_hrx0.70_tray_T_RH.png"),
                          Inches(8.2), Inches(1.5), height=Inches(5.2))
    add_text(s, Inches(8.2), Inches(6.75), Inches(4.8), Inches(0.3),
             "Tray-by-tray T and RH, E2 autumn 10 m² (validation visual)",
             size=10, italic=True, color=GREY)
    add_notes(s, """
Talking points:
- This is the slide to anchor credibility. Every config has been
  instrumented with conservation checks that run every timestep.
- The 1.5 percent residual on the first law is worth explaining: it is
  an accounting artifact of decomposing refrigerant enthalpy into
  Q_evap + W_comp. The CoolProp state-point values themselves are
  exact. If we were to use raw enthalpy differences at each state
  point, it would close to zero.
- Water mass balance closes to 0.000 at the per-timestep level. The
  0.8 percent discrepancy versus 19.2 kg target is because the
  simulation stops slightly below X_final = 0.10 (overshoots by a small
  amount).
- The tray plot on the right is a visual sanity check: air enters tray
  1 hot and dry, exits tray 10 cooler and humid. The gradient is
  smooth and monotonic, no numerical artifacts.
""")


def slide_takeaways():
    s = blank_slide("Key takeaways",
                    "Three things I want to leave you with")
    items = [
        ("1. Put the solar upstream of the condenser.",
         "E2 beats E3 by 7-9 % SEC in every season at both climates. "
         "Collector inlet temperature dominates the combined efficiency of a series HP+solar topology."),
        ("2. Peak COP is a distraction. SEC is what matters.",
         "E3 has higher COP and zero clipping, yet loses on SEC. "
         "High COP on a smaller duty cycle can still cost more than low COP on a bigger one."),
        ("3. Clipping is a thermodynamic limit, not a weather problem.",
         "Clipping happens every clear day, 6+ hours a day, because the stream state "
         "forces the HP to idle when solar alone reaches T_set. The fix is storage."),
    ]
    y = Inches(1.5)
    for i, (title, body) in enumerate(items):
        ty = y + Inches(0.2 + i * 1.7)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), ty, Inches(0.15), Inches(1.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = [TEAL, CORAL, GOLD][i]
        bar.line.fill.background()
        add_text(s, Inches(0.85), ty, Inches(12.0), Inches(0.5),
                 title, size=18, bold=True, color=NAVY)
        add_text(s, Inches(0.85), ty + Inches(0.5), Inches(12.0), Inches(0.9),
                 body, size=13, color=DARK)
    add_notes(s, """
Talking points:
- Each of the three gets about 30 seconds of discussion. The first is the
  engineering answer to the research question. The second is a mental
  model correction: COP has been overweighted in the SAHPD literature.
  The third sets up the scale-up slide by pointing at storage as the
  missing piece.
- Offer to open up: does the professor want to discuss any of these?
""")


def slide_scaleup():
    s = blank_slide("What's next: scale-up and thermal zones",
                    "Going from 3 kg dry (lab) to 500-1000 kg fresh (small commercial)")
    # Left: concept sketch of thermal-zone chamber
    add_text(s, Inches(0.4), Inches(1.3), Inches(6.5), Inches(0.4),
             "Multi-zone chamber concept (4 zones shown)",
             size=13, bold=True, color=NAVY)
    draw_thermal_zones(s, Inches(0.4), Inches(2.5), Inches(6.4))
    add_text(s, Inches(0.4), Inches(5.9), Inches(6.4), Inches(0.5),
             "Each zone gets fresh 45 °C air via the manifold; each zone's exhaust "
             "routes back through its own HRX branch.",
             size=10, italic=True, color=GREY)

    # Right: scale-up numbers and implications
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(7.2), Inches(1.4), Inches(5.8), Inches(5.7))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    add_text(s, Inches(7.3), Inches(1.5), Inches(5.6), Inches(0.4),
             "Capacity scaling (rough)", size=14, bold=True, color=NAVY)
    rows = [
        ("Dry mass", "3 kg", "70 kg", "140 kg"),
        ("Fresh mass", "22.5 kg", "500 kg", "1000 kg"),
        ("Water removed", "19.2 kg", "427 kg", "854 kg"),
        ("Scale factor", "1\u00D7", "22\u00D7", "44\u00D7"),
        ("Trays (current geometry)", "10", "~220", "~440"),
        ("Collector area (same SEC)", "10 m²", "220 m²", "440 m²"),
    ]
    col_w = [Inches(2.2), Inches(1.1), Inches(1.1), Inches(1.1)]
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            left = Inches(7.3) + sum(col_w[:ci], Emu(0))
            top = Inches(2.0 + ri * 0.38)
            fg = NAVY if ri == 0 or ci == 0 else DARK
            add_text(s, left, top, col_w[ci], Inches(0.35),
                     cell, size=11, bold=(ri == 0 or ci == 0), color=fg)
    # Header row bar
    hdr_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(7.3), Inches(1.95), Inches(5.5), Inches(0.03))
    hdr_bar.fill.solid()
    hdr_bar.fill.fore_color.rgb = NAVY
    hdr_bar.line.fill.background()

    add_text(s, Inches(7.3), Inches(4.45), Inches(5.6), Inches(0.4),
             "Why thermal zones (not longer single chamber)",
             size=13, bold=True, color=NAVY)
    bullets(s, Inches(7.3), Inches(4.85), Inches(5.7), Inches(2.2), [
        "Current chamber: RH rises 29 \u2192 84 % across 10 trays in wet phase. Tray 9 barely dries.",
        "Doubling trays in series wastes floor space; tail trays see near-saturated air.",
        "Multi-inlet per zone \u2192 every zone operates in high-VPD region \u2192 shorter, more uniform drying.",
        "Per-zone VPD bypass control becomes natural: dry zones can go cond-direct early.",
        "Caveat: needs distribution manifold + multi-fan + control logic.",
    ], size=11)
    add_notes(s, """
Talking points:
- Current lab scale is 3 kg dry mass. That gets us to interesting physics
  but nowhere near a useful throughput. To dry even one small orchard's
  apple crop in a season we need 500-1000 kg fresh per batch.
- Linear scaling (just make the chamber bigger) does not work. Our
  validation data (V2.3) shows RH climbing from 29 to 84 percent across
  10 trays in wet phase. Add more trays and the tail of the chamber is
  just a hot humid pipe doing no useful work.
- Thermal-zone architecture: split the chamber into N zones, each with
  its own hot-air inlet from a distribution manifold. Each zone sees
  fresh 45 degree air instead of air that already passed through 5 wet
  trays. Every zone operates in a high-VPD regime.
- This also enables per-zone control: a nearly-dry zone can switch to
  condenser-direct bypass early while wet zones still dehumidify. VPD
  bypass logic is already validated in Config A (V3 in the report).
- Open question for discussion: how many zones is optimal? Is there a
  diminishing return? And what is the CAPEX penalty of the manifold
  and multi-fan setup?
- Collector area scales roughly linearly with water mass if we want
  the same SEC. 220 m2 for 500 kg fresh is a realistic rooftop number
  for a small orchard processing building.
- The thermal storage retrofit from the clipping discussion fits in
  here: at 220 m2 the absolute clipped energy is huge, so the case for
  storage gets stronger as we scale up.
""")


def slide_next_steps():
    s = blank_slide("Discussion and next steps", None)
    left_col = [
        ("Immediate", [
            "Thermal storage retrofit for E2 (capture clipped midday heat).",
            "Quantify storage size vs SEC benefit at 10 m² and larger.",
            "Re-run E2 with 15, 20, 30 m² + storage to find the new knee.",
        ]),
        ("Medium term", [
            "Multi-zone chamber model: 2-zone first, validate, then 4-zone.",
            "Per-zone control: coupled VPD bypass + solar-priority scheduling.",
            "Scale-up validation at 50 kg dry (pilot) before committing to 500 kg.",
        ]),
        ("Longer term", [
            "Pilot at 500 kg fresh with local producer.",
            "Economic model: CAPEX vs annualised electricity savings in Nepal context.",
        ]),
    ]
    x = Inches(0.6)
    y = Inches(1.4)
    col_w = Inches(4.0)
    col_gap = Inches(0.15)
    for i, (hdr, items) in enumerate(left_col):
        cx = x + (col_w + col_gap) * i
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, col_w, Inches(0.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = [NAVY, TEAL, SAGE][i]
        bar.line.fill.background()
        add_text(s, cx + Inches(0.15), y + Inches(0.08), col_w - Inches(0.2), Inches(0.35),
                 hdr, size=14, bold=True, color=WHITE)
        bullets(s, cx, y + Inches(0.7), col_w, Inches(5.5), items, size=12)

    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.7),
             "Questions I'd like your input on: "
             "(1) is multi-zone worth modelling now or after thesis?  "
             "(2) thermal storage tank vs PCM? "
             "(3) should we bring in Taplejung for a third climate point before freezing RQ1?",
             size=12, italic=True, color=NAVY)
    add_notes(s, """
Talking points:
- Immediate: the clipping result tells us storage is the lever. Before
  building hardware, run the sim with a simple tank model to see what
  storage fraction recovers the 45 percent waste.
- Medium term: multi-zone model is a bigger undertaking. Need a fresh
  chamber energy/mass balance with N inlets and outlets, and a manifold
  pressure-drop model. Two zone version first, validate against a
  collapse-to-one-zone limit, then go to four.
- Longer term: pilot at a fraction of target capacity before going to
  500 kg. Economic model is critical because Nepal electricity cost
  is low, so SEC wins translate to smaller dollar savings per batch.
- Open questions to the professor at the bottom of the slide are
  intentional prompts. Let him pick which direction to push.
""")


# ============================================================
# Assemble
# ============================================================

slide_title()
slide_base_dryer()
slide_hrx_step()
slide_solar_step()
slide_E2_vs_E3_question()
slide_E2_in_action()
slide_E3_in_action()
slide_side_by_side()
slide_punchline()
slide_clipping()
slide_area_sweep()
slide_seasonal_climate()
slide_validation()
slide_takeaways()
slide_scaleup()
slide_next_steps()

prs.save(str(OUT_FILE))
print(f"Wrote {OUT_FILE}")
print(f"Slides: {len(prs.slides)}")

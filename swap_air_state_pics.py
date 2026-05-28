"""Swap only the air-state pictures on slides 3 & 4. Preserves position, size,
and all other slide content (manual edits stay intact)."""
from pathlib import Path
from pptx import Presentation

DECK = Path(r"C:\Users\sarwa\Downloads\E2_vs_E3_final.pptx")
PLOTS = Path(r"D:\Masters\RQ5\Codex_Chatgpt_drying\RQ1\plots\_audit")

JOBS = [
    (2, "step2e_air_states_E2_kathmandu_annual.png"),
    (3, "step2e_air_states_E3_kathmandu_annual.png"),
]

p = Presentation(DECK)
for slide_idx, png in JOBS:
    slide = p.slides[slide_idx]
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    if len(pics) != 1:
        print(f"slide {slide_idx+1}: expected 1 picture, found {len(pics)} - skipping")
        continue
    pic = pics[0]
    rId = pic._element.blip_rId
    image_part = slide.part.related_part(rId)
    new_bytes = (PLOTS / png).read_bytes()
    image_part._blob = new_bytes
    print(f"slide {slide_idx+1}: swapped image to {png}")

p.save(DECK)
print(f"Saved {DECK}")
